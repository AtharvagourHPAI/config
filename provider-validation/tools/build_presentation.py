"""Generate a branded PowerPoint deck for the Provider Contract Validation Engine.

The deck explains the system end-to-end for a mixed (business + technical)
audience: the problem, the five outcomes, the rulebook, the decision ladder,
the architecture, the data, results, and the Streamlit application.

Run:  python tools/build_presentation.py
Output:  Provider_Contract_Validation_Engine.pptx  (project root)
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
ASSETS = ROOT / "app" / "assets"
OUT = ROOT / "Provider_Contract_Validation_Engine.pptx"

# --- healthplans.ai palette ---
NAVY = RGBColor(0x0B, 0x13, 0x20)
INK = RGBColor(0x16, 0x20, 0x2F)
CYAN = RGBColor(0x5C, 0xC6, 0xE8)
BLUE = RGBColor(0x2F, 0x7D, 0xC6)
SLATE = RGBColor(0x64, 0x73, 0x89)
LIGHT = RGBColor(0xF1, 0xF6, 0xFB)
PANEL = RGBColor(0xFF, 0xFF, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LINE = RGBColor(0xE3, 0xEA, 0xF2)

GREEN = RGBColor(0x1F, 0x9D, 0x6B)
RED = RGBColor(0xD3, 0x49, 0x3F)
AMBER = RGBColor(0xC7, 0x7D, 0x1A)
PURPLE = RGBColor(0x7A, 0x4F, 0xB0)

OUTCOME = {
    "APPROVE": (GREEN, "Apply the change", "All applicable checks passed."),
    "DEVELOP": (BLUE, "Ask for more info", "Recoverable — missing documentation."),
    "DENY": (RED, "No, on the merits", "A substantive disqualifier blocks it."),
    "REJECT": (AMBER, "Return the submission", "Procedurally invalid as filed."),
    "INITIAL_ENROLLMENT_REQUIRED": (
        PURPLE,
        "Route to new enrollment",
        "Out-of-state location add — not a change.",
    ),
}

HEAD = "Segoe UI Semibold"
BODY = "Segoe UI"
MONO = "Consolas"

EMU_W = Inches(13.333)
EMU_H = Inches(7.5)


def _solid(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def rect(slide, x, y, w, h, color, radius=False):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        x, y, w, h,
    )
    _solid(shp, color)
    shp.shadow.inherit = False
    return shp


def text(
    slide, x, y, w, h, runs, size=18, color=INK, bold=False, align=PP_ALIGN.LEFT,
    font=BODY, anchor=MSO_ANCHOR.TOP, line_spacing=1.0,
):
    """Add a textbox. ``runs`` is a string or a list of paragraphs.

    A paragraph may be a string or a list of (text, opts) run-tuples.
    """
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0

    paras = runs if isinstance(runs, list) else [runs]
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        chunks = para if isinstance(para, list) else [(para, {})]
        for chunk_text, opts in chunks:
            r = p.add_run()
            r.text = chunk_text
            r.font.size = Pt(opts.get("size", size))
            r.font.bold = opts.get("bold", bold)
            r.font.name = opts.get("font", font)
            r.font.color.rgb = opts.get("color", color)
    return tb


def bullets(slide, x, y, w, h, items, size=18, gap=8, color=INK):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        lead, rest = (item if isinstance(item, tuple) else (None, item))
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        p.line_spacing = 1.05
        dot = p.add_run()
        dot.text = "\u25AA  "
        dot.font.color.rgb = CYAN
        dot.font.size = Pt(size)
        dot.font.bold = True
        if lead:
            b = p.add_run()
            b.text = lead + "  "
            b.font.bold = True
            b.font.size = Pt(size)
            b.font.name = HEAD
            b.font.color.rgb = INK
        r = p.add_run()
        r.text = rest
        r.font.size = Pt(size)
        r.font.name = BODY
        r.font.color.rgb = color
    return tb


def footer(slide, idx):
    logo = ASSETS / "logo.png"
    if logo.exists():
        slide.shapes.add_picture(
            str(logo), Inches(0.5), Inches(7.04), height=Inches(0.22)
        )
    text(
        slide, Inches(12.0), Inches(7.0), Inches(1.0), Inches(0.3),
        f"{idx:02d}", size=11, color=SLATE, align=PP_ALIGN.RIGHT,
    )


def base_slide(prs, color=WHITE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    rect(slide, 0, 0, EMU_W, EMU_H, color)
    return slide


def content_slide(prs, idx, kicker, title):
    slide = base_slide(prs, WHITE)
    rect(slide, 0, 0, Inches(0.16), EMU_H, CYAN)  # left accent
    text(
        slide, Inches(0.6), Inches(0.42), Inches(11), Inches(0.3),
        kicker.upper(), size=12.5, color=BLUE, bold=True, font=HEAD,
    )
    text(
        slide, Inches(0.58), Inches(0.72), Inches(12), Inches(0.7),
        title, size=30, color=INK, bold=True, font=HEAD,
    )
    rect(slide, Inches(0.6), Inches(1.42), Inches(1.1), Inches(0.05), CYAN)
    footer(slide, idx)
    return slide


# --------------------------------------------------------------------------- #
# Slides
# --------------------------------------------------------------------------- #
def slide_title(prs):
    slide = base_slide(prs, NAVY)
    rect(slide, 0, 0, EMU_W, Inches(0.18), CYAN)
    rect(slide, 0, Inches(7.32), EMU_W, Inches(0.18), CYAN)
    text(
        slide, Inches(0.9), Inches(2.0), Inches(11.5), Inches(0.5),
        "PROVIDER CONFIGURATION · CMS-855I", size=15, color=CYAN, bold=True, font=HEAD,
    )
    text(
        slide, Inches(0.88), Inches(2.5), Inches(11.6), Inches(2.0),
        "Provider Contract Change\nValidation Engine",
        size=48, color=WHITE, bold=True, font=HEAD, line_spacing=1.02,
    )
    text(
        slide, Inches(0.92), Inches(4.45), Inches(11.4), Inches(0.8),
        "A deterministic, fully explainable rule engine that validates provider "
        "contract-change requests and returns one of five auditable outcomes.",
        size=19, color=RGBColor(0xC7, 0xD6, 0xE6), font=BODY, line_spacing=1.15,
    )
    text(
        slide, Inches(0.9), Inches(6.5), Inches(8), Inches(0.4),
        [[("healthplans", {"bold": True, "color": WHITE, "size": 20}),
          (".ai", {"bold": True, "color": CYAN, "size": 20})]],
        font=HEAD,
    )
    return slide


def slide_problem(prs):
    s = content_slide(prs, 2, "The problem", "Why this engine exists")
    bullets(
        s, Inches(0.7), Inches(1.8), Inches(7.2), Inches(5),
        [
            ("Manual & inconsistent.", "Reviewers interpret CMS-855I change requests by hand, so identical cases can get different answers."),
            ("Hard to audit.", "When a decision is questioned, there is no clean record of which rules fired and why."),
            ("Opaque automation.", "An ML classifier would be a black box and can't be trusted for compliance decisions."),
            ("Rare but critical cases.", "The highest-risk outcomes (deny, reject) are the least frequent and easiest to get wrong."),
        ],
        size=17, gap=14,
    )
    # callout panel
    rect(s, Inches(8.3), Inches(1.8), Inches(4.4), Inches(4.4), LIGHT, radius=True)
    text(
        s, Inches(8.65), Inches(2.05), Inches(3.8), Inches(0.5),
        "THE GOAL", size=13, color=BLUE, bold=True, font=HEAD,
    )
    text(
        s, Inches(8.65), Inches(2.5), Inches(3.8), Inches(3.6),
        [
            [("Same inputs \u2192 same answer.", {"bold": True, "size": 17, "color": INK})],
            [("Every decision cites the exact rules, verdicts, and tags that produced it.", {"size": 15, "color": SLATE})],
            [(" ", {"size": 6})],
            [("Rule-driven, not ML-driven.", {"bold": True, "size": 17, "color": INK})],
            [("No randomness, no wall-clock, no hidden state in the decision path.", {"size": 15, "color": SLATE})],
        ],
        line_spacing=1.1,
    )
    return s


def slide_pipeline(prs):
    s = content_slide(prs, 3, "How it works", "One request, one deterministic decision")
    steps = [
        ("Change\nRequest", CYAN, "Inbound CMS-855I edit"),
        ("Current\nContract", CYAN, "Source-of-truth state"),
        ("Tag\nEngine", BLUE, "Derive boolean signals"),
        ("Rulebook\nR-001\u2013R-010", BLUE, "Evaluate applicable rules"),
        ("Decision", NAVY, "Outcome + action + trace"),
    ]
    x = Inches(0.7)
    box_w = Inches(2.15)
    gap = Inches(0.27)
    y = Inches(2.3)
    for i, (label, color, sub) in enumerate(steps):
        rect(s, x, y, box_w, Inches(1.5), color, radius=True)
        text(
            s, x, y + Inches(0.28), box_w, Inches(0.95), label,
            size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER, font=HEAD,
            anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0,
        )
        text(
            s, x - Inches(0.05), y + Inches(1.6), box_w + Inches(0.1), Inches(0.7),
            sub, size=11.5, color=SLATE, align=PP_ALIGN.CENTER,
        )
        if i < len(steps) - 1:
            ar = s.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW, x + box_w + Inches(0.01),
                y + Inches(0.55), gap - Inches(0.02), Inches(0.4),
            )
            _solid(ar, RGBColor(0xB9, 0xCB, 0xDD))
        x = Emu(int(x) + int(box_w) + int(gap))
    text(
        s, Inches(0.7), Inches(5.0), Inches(12), Inches(1.5),
        [
            [("The decision core is a standalone, importable package", {"bold": True, "size": 17, "color": INK}),
             (" \u2014 no UI imports inside ", {"size": 17, "color": INK}),
             ("engine/", {"size": 16, "color": BLUE, "font": MONO}),
             (".  The Streamlit app and the batch scorer both call the same ", {"size": 17, "color": INK}),
             ("decide()", {"size": 16, "color": BLUE, "font": MONO}),
             (".", {"size": 17, "color": INK})],
        ],
        line_spacing=1.15,
    )
    return s


def slide_outcomes(prs):
    s = content_slide(prs, 4, "The vocabulary", "Exactly five possible outcomes")
    order = ["APPROVE", "DEVELOP", "DENY", "REJECT", "INITIAL_ENROLLMENT_REQUIRED"]
    y = Inches(1.85)
    for name in order:
        color, gist, sub = OUTCOME[name]
        rect(s, Inches(0.7), y, Inches(0.16), Inches(0.86), color)
        rect(s, Inches(0.86), y, Inches(11.8), Inches(0.86), LIGHT)
        text(
            s, Inches(1.1), y + Inches(0.12), Inches(4.3), Inches(0.6),
            name.replace("_", " "), size=17, color=color, bold=True, font=HEAD,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        text(
            s, Inches(5.3), y + Inches(0.12), Inches(2.7), Inches(0.6),
            gist, size=15, color=INK, bold=True, anchor=MSO_ANCHOR.MIDDLE,
        )
        text(
            s, Inches(8.0), y + Inches(0.12), Inches(4.4), Inches(0.6),
            sub, size=13.5, color=SLATE, anchor=MSO_ANCHOR.MIDDLE,
        )
        y = Emu(int(y) + int(Inches(1.02)))
    return s


def slide_distinctions(prs):
    s = content_slide(prs, 5, "Telling them apart", "The distinctions that matter")
    # Two panels
    rect(s, Inches(0.7), Inches(1.8), Inches(5.85), Inches(4.7), LIGHT, radius=True)
    text(s, Inches(1.0), Inches(2.0), Inches(5.3), Inches(0.4),
         "APPROVE vs. DEVELOP", size=18, color=BLUE, bold=True, font=HEAD)
    bullets(
        s, Inches(1.0), Inches(2.55), Inches(5.3), Inches(3.8),
        [
            ("APPROVE", "all checks passed \u2014 the change is applied. Final."),
            ("DEVELOP", "might be valid, but information/documentation is missing. A development letter goes out; the case is paused, not closed."),
            ("Key idea", "DEVELOP is recoverable \u2014 \u201cnot yet,\u201d not \u201cno.\u201d"),
        ],
        size=14.5, gap=10,
    )
    rect(s, Inches(6.75), Inches(1.8), Inches(5.85), Inches(4.7), LIGHT, radius=True)
    text(s, Inches(7.05), Inches(2.0), Inches(5.3), Inches(0.4),
         "DENY vs. REJECT", size=18, color=RED, bold=True, font=HEAD)
    bullets(
        s, Inches(7.05), Inches(2.55), Inches(5.3), Inches(3.8),
        [
            ("DENY", "a substantive \u201cno\u201d on the merits (e.g. exclusion/debarment). Referred for review."),
            ("REJECT", "procedurally invalid \u2014 wrong form or unsigned. Returned to the provider to fix and resubmit."),
            ("Key idea", "REJECT never reached the merits; DENY did."),
        ],
        size=14.5, gap=10,
    )
    text(
        s, Inches(0.7), Inches(6.65), Inches(12), Inches(0.5),
        [[("INITIAL_ENROLLMENT_REQUIRED", {"bold": True, "color": PURPLE, "size": 14.5, "font": HEAD}),
          ("  \u2014 re-routes an out-of-state practice-location add to a brand-new enrollment; it is neither a yes nor a no.", {"size": 14.5, "color": INK})]],
    )
    return s


def slide_rules(prs):
    s = content_slide(prs, 6, "The rulebook", "R-001 \u2013 R-010 (editable in YAML)")
    rows = [
        ("Rule", "Passes when\u2026", "On failure"),
        ("R-001", "Form is CMS-855I", "REJECT"),
        ("R-002", "Provider identifiers match the contract", "DEVELOP"),
        ("R-003", "Not an out-of-state practice-location add", "INITIAL ENROLLMENT"),
        ("R-004", "Special payment address is an allowed type", "DEVELOP"),
        ("R-005", "No final adverse legal action (exclusion/debarment)", "DENY"),
        ("R-006", "Reassignee is an enrolled Medicare group", "DEVELOP"),
        ("R-007", "EFT change has CMS-588 + bank docs", "DEVELOP"),
        ("R-008", "Certification statement is signed", "REJECT"),
        ("R-009", "Specialty compatible with provider type", "DENY"),
        ("R-010", "Documentation complete / reusable", "DEVELOP"),
    ]
    tbl = s.shapes.add_table(
        len(rows), 3, Inches(0.7), Inches(1.7), Inches(11.9), Inches(5.1)
    ).table
    tbl.columns[0].width = Inches(1.3)
    tbl.columns[1].width = Inches(7.6)
    tbl.columns[2].width = Inches(3.0)
    for c in range(3):
        cell = tbl.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        p = cell.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = rows[0][c]
        run.font.color.rgb = WHITE
        run.font.bold = True
        run.font.size = Pt(14)
        run.font.name = HEAD
    verdict_color = {"REJECT": AMBER, "DENY": RED, "DEVELOP": BLUE, "INITIAL ENROLLMENT": PURPLE}
    for r in range(1, len(rows)):
        for c in range(3):
            cell = tbl.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if r % 2 else LIGHT
            p = cell.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = rows[r][c]
            run.font.size = Pt(12.5)
            run.font.name = MONO if c == 0 else BODY
            if c == 0:
                run.font.bold = True
                run.font.color.rgb = INK
            elif c == 2:
                run.font.bold = True
                run.font.color.rgb = verdict_color.get(rows[r][c], SLATE)
            else:
                run.font.color.rgb = INK
    return s


def slide_ladder(prs):
    s = content_slide(prs, 7, "Precedence", "The decision ladder \u2014 first match wins")
    tiers = [
        ("1", "DENY", RED, "Adverse exclusion/debarment (R-005) or incompatible specialty (R-009)", "Substantive hard block"),
        ("2", "REJECT", AMBER, "Wrong CMS form (R-001) or unsigned certification (R-008)", "Procedurally invalid"),
        ("3", "INITIAL ENROLLMENT", PURPLE, "Out-of-state practice-location add (R-003)", "Cannot be a modification"),
        ("4", "DEVELOP", BLUE, "Any DEVELOP-class rule (R-002/004/006/007/010) or docs incomplete", "Needs more information"),
        ("5", "APPROVE", GREEN, "No rule blocked the request", "Default \u2014 apply the change"),
    ]
    y = Inches(1.8)
    for num, name, color, cond, note in tiers:
        rect(s, Inches(0.7), y, Inches(0.8), Inches(0.86), color, radius=True)
        text(s, Inches(0.7), y, Inches(0.8), Inches(0.86), num, size=26, color=WHITE,
             bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=HEAD)
        rect(s, Inches(1.6), y, Inches(11.05), Inches(0.86), LIGHT)
        text(s, Inches(1.85), y + Inches(0.1), Inches(3.4), Inches(0.66), name,
             size=16, color=color, bold=True, font=HEAD, anchor=MSO_ANCHOR.MIDDLE)
        text(s, Inches(5.0), y + Inches(0.06), Inches(7.4), Inches(0.74),
             [[(cond, {"size": 13.5, "color": INK})],
              [(note, {"size": 11.5, "color": SLATE})]],
             anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
        y = Emu(int(y) + int(Inches(1.0)))
    return s


def slide_data(prs):
    s = content_slide(prs, 8, "The data", "Five reference sources")
    rows = [
        ("Existing_Provider_Contracts", "80", "Source-of-truth contract state"),
        ("Requested_Contract_Changes", "60 + 14", "Labeled change requests (+ synthetic review rows)"),
        ("Validation_Rules", "10", "Human-readable rulebook"),
        ("Lists", "5", "Controlled outcome vocabulary"),
        ("Summary", "\u2014", "Expected counts / metadata"),
    ]
    y = Inches(1.8)
    for name, vol, desc in rows:
        rect(s, Inches(0.7), y, Inches(4.7), Inches(0.7), LIGHT)
        text(s, Inches(0.9), y + Inches(0.08), Inches(4.4), Inches(0.55), name,
             size=13.5, color=INK, bold=True, font=MONO, anchor=MSO_ANCHOR.MIDDLE)
        rect(s, Inches(5.45), y, Inches(1.1), Inches(0.7), NAVY)
        text(s, Inches(5.45), y + Inches(0.08), Inches(1.1), Inches(0.55), vol,
             size=15, color=CYAN, bold=True, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE, font=HEAD)
        text(s, Inches(6.75), y + Inches(0.08), Inches(5.9), Inches(0.55), desc,
             size=13.5, color=SLATE, anchor=MSO_ANCHOR.MIDDLE)
        y = Emu(int(y) + int(Inches(0.82)))
    text(
        s, Inches(0.7), Inches(6.0), Inches(12), Inches(1.2),
        [
            [("Design facts baked in:", {"bold": True, "size": 14.5, "color": INK})],
            [("0 orphan links (every change resolves to a real contract) \u00b7 labels are heavily imbalanced \u00b7 ", {"size": 13.5, "color": SLATE}),
             ("REJECT", {"size": 13.5, "color": AMBER, "bold": True}),
             (" is absent from real data, so it is produced by rules and exercised with synthetic rows.", {"size": 13.5, "color": SLATE})],
        ],
        line_spacing=1.1,
    )
    return s


def slide_results(prs):
    s = content_slide(prs, 9, "Results", "100% reproduction of labeled outcomes")
    # Big accuracy stat
    rect(s, Inches(0.7), Inches(1.9), Inches(3.6), Inches(2.0), NAVY, radius=True)
    text(s, Inches(0.7), Inches(2.1), Inches(3.6), Inches(1.0), "74 / 74",
         size=44, color=WHITE, bold=True, align=PP_ALIGN.CENTER, font=HEAD)
    text(s, Inches(0.7), Inches(3.15), Inches(3.6), Inches(0.6),
         "outcomes matched (100%)", size=15, color=CYAN, align=PP_ALIGN.CENTER, font=HEAD)
    # Distribution bars
    dist = [("APPROVE", 49, GREEN), ("DEVELOP", 12, BLUE), ("DENY", 10, RED),
            ("REJECT", 2, AMBER), ("INITIAL_ENROLLMENT_REQUIRED", 1, PURPLE)]
    y = Inches(1.9)
    max_w = Inches(5.2)
    for name, n, color in dist:
        text(s, Inches(4.7), y, Inches(3.4), Inches(0.4), name.replace("_", " "),
             size=12.5, color=INK, bold=True, anchor=MSO_ANCHOR.MIDDLE, font=HEAD)
        bar_w = Emu(int(Inches(0.35)) + int(max_w * (n / 49)))
        rect(s, Inches(7.0), y + Inches(0.03), bar_w, Inches(0.34), color, radius=True)
        text(s, Emu(int(Inches(7.0)) + int(bar_w) + int(Inches(0.1))), y, Inches(0.8),
             Inches(0.4), str(n), size=13, color=INK, bold=True, anchor=MSO_ANCHOR.MIDDLE)
        y = Emu(int(y) + int(Inches(0.92)))
    text(
        s, Inches(0.7), Inches(6.3), Inches(12), Inches(0.8),
        [[("Every DENY traces to R-005/R-009, the single INITIAL_ENROLLMENT to R-003, "
           "every REJECT to R-001/R-008 \u2014 asserted by the regression test suite (30 tests).",
           {"size": 14, "color": SLATE})]],
        line_spacing=1.15,
    )
    return s


def slide_architecture(prs):
    s = content_slide(prs, 10, "Architecture", "A pure engine, a thin UI, editable config")
    cols = [
        ("engine/  (decision core)", BLUE, [
            "models.py \u00b7 enums.py \u2014 typed vocabulary",
            "loaders.py \u2014 sheets \u2192 models + integrity",
            "tag_engine.py \u2014 derive signals",
            "rules.py \u00b7 rule_engine.py \u2014 R-001\u2013R-010",
            "decision_engine.py \u2014 precedence ladder",
            "scoring.py \u2014 batch metrics",
        ]),
        ("app/  (Streamlit UI)", NAVY, [
            "input_page \u2014 upload & process",
            "outcome_explorer \u2014 cards drill-down",
            "selector \u00b7 trace_viewer \u2014 single decision",
            "batch_report \u2014 accuracy & confusion",
            "reference_viewer \u2014 source data",
            "theme \u2014 healthplans.ai styling",
        ]),
        ("config/  (no-code tuning)", GREEN, [
            "rules.yaml \u2014 rule scope, verdicts,",
            "   and matching vocabularies",
            "outcomes.yaml \u2014 outcome list +",
            "   default contract actions",
            "",
            "Retune without touching code.",
        ]),
    ]
    x = Inches(0.7)
    w = Inches(3.92)
    for title_, color, items in cols:
        rect(s, x, Inches(1.8), w, Inches(0.6), color)
        text(s, x + Inches(0.15), Inches(1.9), w - Inches(0.3), Inches(0.4), title_,
             size=14.5, color=WHITE, bold=True, font=HEAD, anchor=MSO_ANCHOR.MIDDLE)
        rect(s, x, Inches(2.4), w, Inches(4.0), LIGHT)
        tb = s.shapes.add_textbox(x + Inches(0.2), Inches(2.6), w - Inches(0.4), Inches(3.7))
        tf = tb.text_frame
        tf.word_wrap = True
        for i, it in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_after = Pt(7)
            r = p.add_run()
            r.text = it
            r.font.size = Pt(12.5)
            r.font.name = MONO if "." in it or it.startswith(" ") else BODY
            r.font.color.rgb = INK if it and not it.startswith(" ") else SLATE
        x = Emu(int(x) + int(w) + int(Inches(0.17)))
    return s


def slide_app(prs):
    s = content_slide(prs, 11, "The application", "Five views for reviewers")
    cards = [
        ("Input", CYAN, "Upload a change-request workbook (or use the sample) and process it."),
        ("Outcomes", BLUE, "Clickable result cards \u2192 contract list \u2192 plain-language explanation."),
        ("Single decision", NAVY, "Pick a request; see change vs. contract and the full rule trace."),
        ("Batch report", GREEN, "Accuracy, expected-vs-actual distribution, confusion matrix, mismatches."),
        ("Reference data", PURPLE, "Browse the source-of-truth sheets with summaries from the sidebar."),
    ]
    positions = [
        (Inches(0.7), Inches(1.85)), (Inches(4.75), Inches(1.85)), (Inches(8.8), Inches(1.85)),
        (Inches(0.7), Inches(4.25)), (Inches(4.75), Inches(4.25)),
    ]
    cw, ch = Inches(3.85), Inches(2.15)
    for (title_, color, desc), (x, y) in zip(cards, positions):
        rect(s, x, y, cw, ch, LIGHT, radius=True)
        rect(s, x, y, cw, Inches(0.12), color)
        text(s, x + Inches(0.25), y + Inches(0.3), cw - Inches(0.5), Inches(0.5), title_,
             size=18, color=color, bold=True, font=HEAD)
        text(s, x + Inches(0.25), y + Inches(0.85), cw - Inches(0.5), Inches(1.1), desc,
             size=13.5, color=SLATE, line_spacing=1.12)
    # highlight new feature
    rect(s, Inches(8.8), Inches(4.25), Inches(3.85), Inches(2.15), NAVY, radius=True)
    text(s, Inches(9.05), Inches(4.55), Inches(3.4), Inches(0.5), "NEW", size=13,
         color=CYAN, bold=True, font=HEAD)
    text(s, Inches(9.05), Inches(4.95), Inches(3.4), Inches(1.4),
         "Outcome cards drill down to a layman summary explaining what the engine "
         "decided and why \u2014 with the technical trace one click away.",
         size=13, color=WHITE, line_spacing=1.12)
    return s


def slide_principles(prs):
    s = content_slide(prs, 12, "Principles", "What makes it trustworthy")
    left = [
        ("Deterministic.", "Identical inputs always yield the identical outcome and trace."),
        ("Explainable.", "Every decision lists the rules evaluated, their verdicts, and the tags that fired."),
        ("Separation of concerns.", "Logic lives in one place; precedence in another; no UI code in the engine."),
    ]
    right = [
        ("Label-blind.", "decide() never reads the ground-truth label fields."),
        ("Editable rulebook.", "Analysts retune scope and verdicts in YAML without code changes."),
        ("Tested.", "30 tests: loaders, each rule, precedence, full regression, synthetic REJECT."),
    ]
    bullets(s, Inches(0.7), Inches(1.95), Inches(5.9), Inches(5), left, size=15.5, gap=16)
    bullets(s, Inches(6.85), Inches(1.95), Inches(5.8), Inches(5), right, size=15.5, gap=16)
    return s


def slide_summary(prs):
    s = base_slide(prs, NAVY)
    rect(s, 0, 0, Inches(0.16), EMU_H, CYAN)
    text(s, Inches(0.9), Inches(0.9), Inches(11), Inches(0.4),
         "IN ONE LINE", size=14, color=CYAN, bold=True, font=HEAD)
    text(
        s, Inches(0.88), Inches(1.4), Inches(11.6), Inches(1.8),
        "A rule-driven engine that turns any CMS-855I change request into one of "
        "five auditable outcomes \u2014 with a complete, human-readable reason every time.",
        size=27, color=WHITE, bold=True, font=HEAD, line_spacing=1.08,
    )
    chips = ["APPROVE", "DEVELOP", "DENY", "REJECT", "INITIAL_ENROLLMENT_REQUIRED"]
    x = Inches(0.9)
    y = Inches(3.7)
    for name in chips:
        color = OUTCOME[name][0]
        w = Inches(1.7 if name != "INITIAL_ENROLLMENT_REQUIRED" else 3.7)
        rect(s, x, y, w, Inches(0.6), color, radius=True)
        text(s, x, y, w, Inches(0.6), name.replace("_", " "), size=13, color=WHITE,
             bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=HEAD)
        x = Emu(int(x) + int(w) + int(Inches(0.2)))
    text(
        s, Inches(0.9), Inches(4.9), Inches(11.5), Inches(1.4),
        [
            [("Deterministic", {"color": CYAN, "bold": True, "size": 18, "font": HEAD}),
             ("   \u00b7   ", {"color": SLATE, "size": 18}),
             ("Explainable", {"color": CYAN, "bold": True, "size": 18, "font": HEAD}),
             ("   \u00b7   ", {"color": SLATE, "size": 18}),
             ("Editable", {"color": CYAN, "bold": True, "size": 18, "font": HEAD}),
             ("   \u00b7   ", {"color": SLATE, "size": 18}),
             ("100% reproduction on labeled data", {"color": CYAN, "bold": True, "size": 18, "font": HEAD})],
        ],
    )
    text(s, Inches(0.9), Inches(6.6), Inches(8), Inches(0.4),
         [[("healthplans", {"bold": True, "color": WHITE, "size": 18}),
           (".ai", {"bold": True, "color": CYAN, "size": 18})]], font=HEAD)
    return s


def build():
    prs = Presentation()
    prs.slide_width = EMU_W
    prs.slide_height = EMU_H

    slide_title(prs)
    slide_problem(prs)
    slide_pipeline(prs)
    slide_outcomes(prs)
    slide_distinctions(prs)
    slide_rules(prs)
    slide_ladder(prs)
    slide_data(prs)
    slide_results(prs)
    slide_architecture(prs)
    slide_app(prs)
    slide_principles(prs)
    slide_summary(prs)

    prs.save(OUT)
    print(f"Wrote {OUT.name}: {len(prs.slides._sldIdLst)} slides -> {OUT}")


if __name__ == "__main__":
    build()
